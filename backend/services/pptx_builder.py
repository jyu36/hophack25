# backend/services/pptx_builder.py
from __future__ import annotations

from typing import Dict, Any, List
from pathlib import Path
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

import matplotlib
matplotlib.use("Agg")  # headless
import matplotlib.pyplot as plt
import numpy as np

# ----------------- Helper: theme -----------------
def _theme_colors(plan: Dict[str, Any]):
    theme = (plan.get("SLIDE_PLAN") or {}).get("theme", {})
    accent_hex = (theme.get("accent") or "#2396F5").lstrip("#")
    try:
        r = int(accent_hex[0:2], 16)
        g = int(accent_hex[2:4], 16)
        b = int(accent_hex[4:6], 16)
    except Exception:
        r, g, b = (35, 150, 245)
    return RGBColor(r, g, b)

# ----------------- Chart renderers -----------------
def _save_plot(fig) -> str:
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
    fig.savefig(tmp.name, bbox_inches="tight", dpi=240)
    plt.close(fig)
    return tmp.name

def _chart_bar(spec: Dict[str, Any]) -> str:
    x = spec.get("x", [])
    y = spec.get("y", [])
    fig, ax = plt.subplots(figsize=(8, 4.5))
    ax.bar(x, y)
    ax.set_xlabel(spec.get("xLabel",""))
    ax.set_ylabel(spec.get("yLabel",""))
    ax.grid(True, axis="y", alpha=0.3)
    return _save_plot(fig)

def _chart_line(spec: Dict[str, Any]) -> str:
    x = spec.get("x", [])
    y = spec.get("y", [])
    fig, ax = plt.subplots(figsize=(8, 4.5))
    ax.plot(x, y, marker="o")
    ax.set_xlabel(spec.get("xLabel",""))
    ax.set_ylabel(spec.get("yLabel",""))
    ax.grid(True, alpha=0.3)
    return _save_plot(fig)

def _chart_scatter_2d(spec: Dict[str, Any]) -> str:
    pts = spec.get("points", [])
    fig, ax = plt.subplots(figsize=(8, 4.5))
    for p in pts:
        ax.scatter(p.get("x",0), p.get("y",0), s=max(10, float(p.get("size", 10))), alpha=0.8)
    ax.set_xlabel(spec.get("xLabel","UMAP-1"))
    ax.set_ylabel(spec.get("yLabel","UMAP-2"))
    ax.grid(True, alpha=0.3)
    return _save_plot(fig)

def _chart_heatmap(spec: Dict[str, Any]) -> str:
    xcats = spec.get("xCats", [])
    ycats = spec.get("yCats", [])
    vals = spec.get("values", [])
    arr = np.array(vals) if len(vals)>0 else np.zeros((max(1,len(ycats)), max(1,len(xcats))))
    fig, ax = plt.subplots(figsize=(8,4.5))
    im = ax.imshow(arr, aspect="auto")
    ax.set_xticks(range(len(xcats)), labels=xcats, rotation=45, ha="right")
    ax.set_yticks(range(len(ycats)), labels=ycats)
    fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    return _save_plot(fig)

def _chart_radar(spec: Dict[str, Any]) -> str:
    indicators = spec.get("indicators", [])
    series = spec.get("series", [])
    N = len(indicators) or 3
    angles = np.linspace(0, 2*np.pi, N, endpoint=False).tolist()
    angles += angles[:1]

    fig = plt.figure(figsize=(6,6))
    ax = plt.subplot(111, polar=True)
    for s in series:
        vals = s.get("values", [0]*N)
        vals = (vals + vals[:1]) if len(vals)==N else ([0]*N+[0])
        ax.plot(angles, vals, linewidth=2)
        ax.fill(angles, vals, alpha=0.1)
    ax.set_xticks(angles[:-1], indicators)
    ax.set_yticklabels([])
    return _save_plot(fig)

def _chart_gantt(spec: Dict[str, Any]) -> str:
    tasks = spec.get("tasks", [])
    fig, ax = plt.subplots(figsize=(8, 4.5))
    for i, t in enumerate(tasks):
        # Placeholder durations: we do not parse dates here; we just plot sequential bars
        ax.barh(i, 1.0, left=i*0.1 + 0.1)
        ax.text(0.05, i, t.get("name", f"T{i}"), va="center", ha="left", transform=ax.get_yaxis_transform())
    ax.set_yticks(range(len(tasks)), labels=[t.get("name","") for t in tasks])
    ax.invert_yaxis()
    ax.set_xlabel("Timeline (placeholder)")
    return _save_plot(fig)

def _chart_flow(spec: Dict[str, Any]) -> str:
    # Minimal linear flow renderer (no extra deps): draw boxes left-to-right with arrows
    nodes = spec.get("nodes", [])
    edges = spec.get("edges", [])
    labels = [n.get("label") or n.get("id") for n in nodes] or ["Ingest","Process","Output"]

    fig, ax = plt.subplots(figsize=(8, 2.8))
    ax.axis("off")
    n = len(labels)
    xs = np.linspace(0.1, 0.9, n)
    y = 0.5
    for i, (x, lab) in enumerate(zip(xs, labels)):
        ax.add_patch(plt.Rectangle((x-0.06, y-0.08), 0.12, 0.16, edgecolor="black", facecolor="white"))
        ax.text(x, y, lab, ha="center", va="center")
        if i < n-1:
            ax.annotate("", xy=(xs[i+1]-0.12, y), xytext=(x+0.06, y),
                        arrowprops=dict(arrowstyle="->", lw=1.5))
    return _save_plot(fig)

def _render_chart_image(chart: Dict[str, Any]) -> str:
    try:
        ctype = (chart.get("type") or "").lower()
        spec = chart.get("spec", {})
        if ctype == "bar":          return _chart_bar(spec)
        if ctype == "line":         return _chart_line(spec)
        if ctype == "scatter-2d":   return _chart_scatter_2d(spec)
        if ctype == "heatmap":      return _chart_heatmap(spec)
        if ctype == "radar":        return _chart_radar(spec)
        if ctype == "gantt":        return _chart_gantt(spec)
        if ctype == "flow-diagram": return _chart_flow(spec)
        # Fallback empty plot
        fig, ax = plt.subplots(figsize=(6,3))
        ax.text(0.5, 0.5, f"Unsupported chart: {ctype}", ha="center", va="center")
        return _save_plot(fig)
    except Exception as e:
        fig, ax = plt.subplots(figsize=(6,3))
        ax.text(0.5, 0.5, f"Chart error: {e}", ha="center", va="center", wrap=True)
        return _save_plot(fig)

# ----------------- PPTX builder -----------------
def _add_title(slide, text: str, color: RGBColor):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = box.text_frame
    tf.clear()
    # set text robustly even if no runs exist yet
    p = tf.paragraphs[0]
    p.text = text or "Untitled"
    p.alignment = PP_ALIGN.LEFT
    # apply font on the first run if present, else create one
    if p.runs:
        run = p.runs[0]
    else:
        run = p.add_run()
        run.text = p.text
    run.font.bold = True
    run.font.size = Pt(36)
    run.font.color.rgb = color

def _add_bullets(slide, bullets: List[str], top=1.5, left=0.6, width=5.2, height=4.5):
    if not bullets:
        return
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = str(b) if b is not None else ""
        # ensure a run exists to set font size
        if p.runs:
            r0 = p.runs[0]
        else:
            r0 = p.add_run(); r0.text = p.text
        r0.font.size = Pt(18)
        p.level = 0

def _add_table(slide, table_spec: Dict[str, Any], left=0.6, top=1.5, width=4.8, height=4.0):
    cols = table_spec.get("columns", [])
    rows = table_spec.get("rows", [])
    if not cols: return
    rows_count = max(1, len(rows)+1)
    shape = slide.shapes.add_table(rows_count, len(cols),
                                   Inches(left), Inches(top), Inches(width), Inches(height)).table
    # header
    for j, col in enumerate(cols):
        shape.cell(0, j).text = str(col)
    # rows
    for i, r in enumerate(rows, start=1):
        for j, cell in enumerate(r[:len(cols)]):
            shape.cell(i, j).text = str(cell)

def _add_image(slide, img_path: str, left=6.0, top=1.5, width=4.0):
    try:
        slide.shapes.add_picture(img_path, Inches(left), Inches(top), width=Inches(width))
    except Exception:
        pass

def build_pptx(plan: Dict[str, Any], out_path: str) -> str:
    """
    Build a .pptx from the {SLIDE_PLAN, SPEAKER_NOTES} dict produced by the LLM (or stub).
    We render charts via matplotlib to PNGs, then place them on slides.
    """
    prs = Presentation()
    accent = _theme_colors(plan)

    slides = (plan.get("SLIDE_PLAN") or {}).get("slides", [])
    if not isinstance(slides, list):
        slides = []
    if not slides:
        # Add a minimal placeholder slide to avoid empty deck
        s = prs.slides.add_slide(prs.slide_layouts[5])
        _add_title(s, "No content available", accent)
        _add_bullets(s, ["The generated plan contained no slides."])
        Path(out_path).parent.mkdir(exist_ok=True, parents=True)
        prs.save(out_path)
        return out_path
    for slide_def in slides:
        layout = slide_def.get("layout", "title-only")
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

        _add_title(slide, slide_def.get("title", ""), accent)

        bullets = slide_def.get("bullets") or []
        tables = slide_def.get("tables") or []
        charts = slide_def.get("charts") or []

        if layout == "title-only":
            if not bullets:
                bullets = ["(placeholder) Add key highlights here"]
            _add_bullets(slide, bullets, top=1.8, left=1.0, width=8.5, height=4.5)

        elif layout == "two-col":
            if not bullets:
                bullets = ["(placeholder) Add bullet points"]
            _add_bullets(slide, bullets, top=1.5, left=0.6, width=4.8, height=4.5)
            # chart or table on right
            if charts:
                try:
                    img = _render_chart_image(charts[0])
                    _add_image(slide, img, left=6.0, top=1.5, width=4.0)
                except Exception:
                    pass
            elif tables:
                _add_table(slide, tables[0], left=6.0, top=1.5, width=4.0, height=4.0)

        elif layout == "table-left-chart-right":
            if tables:
                _add_table(slide, tables[0], left=0.6, top=1.5, width=4.8, height=4.0)
            else:
                if not bullets:
                    bullets = ["(placeholder) Add details"]
                _add_bullets(slide, bullets, top=1.5, left=0.6, width=4.8, height=4.5)
            if charts:
                try:
                    img = _render_chart_image(charts[0])
                    _add_image(slide, img, left=6.0, top=1.5, width=4.0)
                except Exception:
                    pass

        elif layout == "chart-full":
            if charts:
                try:
                    img = _render_chart_image(charts[0])
                    _add_image(slide, img, left=1.2, top=1.5, width=8.0)
                except Exception:
                    pass
            if not bullets:
                bullets = ["(placeholder) Chart caption"]
            _add_bullets(slide, bullets, top=5.2, left=1.2, width=8.0, height=1.5)

        else:
            # default fallback
            if not bullets:
                bullets = ["(placeholder) Content"]
            _add_bullets(slide, bullets, top=1.8, left=1.0, width=8.5, height=4.5)

    Path(out_path).parent.mkdir(exist_ok=True, parents=True)
    prs.save(out_path)
    return out_path
